import io
from .base import BaseParser

class PPTXParser(BaseParser):
    def parse(self, file_content: bytes) -> dict:
        text = self.extract_text(file_content)
        meta = self.extract_metadata(file_content)
        tables = self.extract_tables(file_content)
        images = self.extract_images(file_content)
        structure = self.extract_structure(file_content)

        words = len(text.split())
        chars = len(text)
        
        return {
            "title": meta.get("title", ""),
            "metadata": meta,
            "sections": structure,
            "paragraphs": [{"text": p.strip()} for p in text.split("\n\n") if p.strip()],
            "tables": tables,
            "images": images,
            "links": [],
            "language": "en",
            "statistics": {
                "word_count": words,
                "character_count": chars,
                "page_count": meta.get("slide_count", 1),
                "table_count": len(tables),
                "image_count": len(images)
            }
        }

    def extract_text(self, file_content: bytes) -> str:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_content))
            slides_text = []
            for slide in prs.slides:
                slide_paragraphs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_paragraphs.append(shape.text.strip())
                slides_text.append("\n".join(slide_paragraphs))
            return "\n\n".join(slides_text)
        except ImportError:
            return "python-pptx not installed. Return empty fallback."
        except Exception as e:
            return f"Error parsing PPTX: {str(e)}"

    def extract_tables(self, file_content: bytes) -> list[dict]:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_content))
            tables_list = []
            for s_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        grid = []
                        for row in table.rows:
                            grid.append([cell.text.strip() for cell in row.cells])
                        tables_list.append({
                            "slide": s_idx + 1,
                            "data": grid
                        })
            return tables_list
        except Exception:
            return []

    def extract_images(self, file_content: bytes) -> list[dict]:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_content))
            images_list = []
            for s_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13: # Picture shape type
                        images_list.append({
                            "slide": s_idx + 1,
                            "name": shape.name,
                            "width": shape.width,
                            "height": shape.height
                        })
            return images_list
        except Exception:
            return []

    def extract_metadata(self, file_content: bytes) -> dict:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_content))
            return {
                "slide_count": len(prs.slides),
                "title": prs.core_properties.title or ""
            }
        except Exception:
            return {"slide_count": 1}

    def extract_structure(self, file_content: bytes) -> list[dict]:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_content))
            structure = []
            for s_idx, slide in enumerate(prs.slides):
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {s_idx + 1}"
                structure.append({
                    "level": "1",
                    "title": title,
                    "slide": s_idx + 1
                })
            return structure
        except Exception:
            return []
